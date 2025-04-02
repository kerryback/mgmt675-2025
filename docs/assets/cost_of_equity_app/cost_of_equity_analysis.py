import pandas_datareader.data as web
import yfinance as yf
import pandas as pd
import numpy as np
from datetime import datetime, timedelta
import matplotlib.pyplot as plt
from scipy import stats
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import io
from lxml import etree

def get_market_risk_premium():
    """Calculate market risk premium using Fama-French data"""
    ff_data = web.DataReader('F-F_Research_Data_Factors', 'famafrench', start='1926-07-01')[0]
    monthly_mrp = ff_data['Mkt-RF'].mean()
    annual_mrp_decimal = (monthly_mrp * 12) / 100
    return annual_mrp_decimal

def calculate_beta(ticker):
    """Calculate beta and create regression plot for a given ticker"""
    end_date = datetime.now()
    start_date = end_date - timedelta(days=3650)
    
    ff_data = web.DataReader('F-F_Research_Data_Factors', 'famafrench', 
                            start=start_date.strftime('%Y-%m-%d'))[0]
    
    ff_data['Mkt-RF'] = ff_data['Mkt-RF'] / 100
    ff_data['RF'] = ff_data['RF'] / 100
    
    stock_data = yf.download(ticker, start=start_date, end=end_date, interval='1mo')
    stock_returns = stock_data['Close'].pct_change()
    
    stock_returns.index = pd.PeriodIndex(stock_returns.index.year * 100 + stock_returns.index.month, freq='M')
    ff_data.index = pd.PeriodIndex(ff_data.index.year * 100 + ff_data.index.month, freq='M')
    
    merged_data = pd.concat([stock_returns, ff_data['Mkt-RF'], ff_data['RF']], axis=1)
    merged_data.columns = ['Stock_Return', 'Mkt_RF', 'RF']
    
    merged_data = merged_data.dropna()
    merged_data = merged_data.tail(60)
    
    merged_data['Excess_Return'] = merged_data['Stock_Return'] - merged_data['RF']
    
    slope, intercept, r_value, p_value, std_err = stats.linregress(
        merged_data['Mkt_RF'], 
        merged_data['Excess_Return']
    )
    
    plt.figure(figsize=(10, 6))
    plt.scatter(merged_data['Mkt_RF'], merged_data['Excess_Return'], alpha=0.5)
    plt.plot(merged_data['Mkt_RF'], intercept + slope * merged_data['Mkt_RF'], 'r')
    plt.xlabel('Market Risk Premium (Mkt-RF)')
    plt.ylabel('Stock Excess Return')
    plt.title(f'Beta Calculation for {ticker}\nBeta = {slope:.2f}')
    plt.grid(True, alpha=0.3)
    
    return slope, plt.gcf()

def get_risk_free_rate():
    """Get current risk-free rate from FRED (3-month T-bill)"""
    end_date = datetime.now()
    start_date = end_date - timedelta(days=60)
    
    tb3ms = web.DataReader('TB3MS', 'fred', start_date, end_date)
    tb3ms = tb3ms.dropna()
    rf_rate = tb3ms['TB3MS'].iloc[-1] / 100
    
    return rf_rate

def parse_xml(xml_string):
    """Helper function to parse XML for PowerPoint formatting"""
    return etree.fromstring(xml_string.encode('utf-8'))

def create_cost_of_equity_deck(ticker):
    """
    Create a PowerPoint presentation analyzing the cost of equity for a given stock ticker.
    
    Parameters:
    ticker (str): Stock ticker symbol (e.g., 'AAPL', 'MSFT', 'WMT')
    
    Returns:
    str: Filename of the saved PowerPoint presentation
    dict: Results dictionary containing the calculated values
    """
    try:
        # Get cost of equity calculations
        mrp = get_market_risk_premium()
        beta, plot = calculate_beta(ticker)
        rf = get_risk_free_rate()
        cost_of_equity = rf + beta * mrp
        
        results = {
            'market_risk_premium': mrp,
            'beta': beta,
            'risk_free_rate': rf,
            'cost_of_equity': cost_of_equity,
            'plot': plot
        }
        
        # Create presentation
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text = f"Cost of Equity Analysis for {ticker.upper()}"
        
        # Beta plot slide
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = f"Beta Calculation for {ticker.upper()}"
        
        # Save the plot to a bytes buffer
        buf = io.BytesIO()
        results['plot'].savefig(buf, format='png', bbox_inches='tight', dpi=300)
        buf.seek(0)
        
        # Add the plot to the slide
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(4.5)
        slide.shapes.add_picture(buf, left, top, width, height)
        
        # Results table slide
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Cost of Equity Components"
        
        # Create table
        rows = 5
        cols = 3
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(4)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Set column headers
        table.cell(0, 0).text = "Component"
        table.cell(0, 1).text = "Value"
        table.cell(0, 2).text = "Source"
        
        # Fill table data
        data = [
            ["Risk-free Rate", f"{results['risk_free_rate']:.2%}", "3-month T-bill (FRED)"],
            ["Beta", f"{results['beta']:.2f}", "Regression using monthly returns"],
            ["Market Risk Premium", f"{results['market_risk_premium']:.2%}", "Fama-French factors (1926-present)"],
            ["Cost of Equity", f"{results['cost_of_equity']:.2%}", "rf + β × MRP"]
        ]
        
        # Fill table and format
        for row_idx, (item, value, source) in enumerate(data, start=1):
            table.cell(row_idx, 0).text = item
            table.cell(row_idx, 1).text = value
            table.cell(row_idx, 2).text = source
            
            # Add alternating row colors
            for col in range(3):
                cell = table.cell(row_idx, col)
                if row_idx % 2 == 1:
                    cell._tc.get_or_add_tcPr().append(parse_xml(
                        '<w:shd xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main" w:fill="E8E8E8"/>'
                    ))
        
        # Format column headers
        for col in range(3):
            cell = table.cell(0, col)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.bold = True
            cell._tc.get_or_add_tcPr().append(parse_xml(
                '<w:shd xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main" w:fill="4472C4"/>'
            ))
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
        
        # Save presentation
        filename = f"{ticker.upper()}_cost_of_equity.pptx"
        prs.save(filename)
        
        # Print the results
        print(f"\nCost of Equity Analysis Results for {ticker.upper()}")
        print(f"Risk-free Rate: {results['risk_free_rate']:.2%}")
        print(f"Beta: {results['beta']:.2f}")
        print(f"Market Risk Premium: {results['market_risk_premium']:.2%}")
        print(f"Cost of Equity: {results['cost_of_equity']:.2%}")
        
        return filename, results
        
    except Exception as e:
        raise Exception(f"Error creating presentation for {ticker}: {str(e)}")

if __name__ == "__main__":
    # Example usage
    ticker = "AAPL"  # or any other ticker
    try:
        filename, results = create_cost_of_equity_deck(ticker)
        print(f"\nPresentation saved as: {filename}")
        plt.show()  # Display the beta plot
    except Exception as e:
        print(f"Error: {str(e)}")
