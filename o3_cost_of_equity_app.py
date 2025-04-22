"""
Streamlit web app that takes a ticker and generates a three‑slide PowerPoint
containing:
  • Title slide with “{TICKER} Cost of Equity Capital”
  • Scatter plot of the stock’s excess return vs. the market excess return with a
    regression line (no confidence interval)
  • Table summarising the CAPM cost‑of‑equity calculation

Dependencies (add these to requirements.txt):
streamlit>=1.32
pandas>=2.2
numpy>=1.26
pandas_datareader>=0.10
statsmodels>=0.14
yfinance==0.2.54
matplotlib>=3.8
python-pptx>=0.6
"""

import datetime as dt
from io import BytesIO

import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
import statsmodels.api as sm
import streamlit as st
from pandas_datareader import data as web
from pptx import Presentation
from pptx.util import Inches
import yfinance as yf


# ---------------------------------------------------------------------------
# Helper functions
# ---------------------------------------------------------------------------

def fetch_market_factors():
    """Download Fama‑French 3 factors (monthly)."""
    ff = web.DataReader("F-F_Research_Data_Factors", "famafrench")[0]
    ff.index = ff.index.to_timestamp("M")  # period → month‑end timestamp
    return ff[["Mkt-RF"]]


def fetch_risk_free(start, end):
    """Download 3‑month T‑bill rate from FRED and resample to month‑end (%)."""
    rf_daily = web.DataReader("DTB3", "fred", start, end)
    rf_monthly = rf_daily.resample("M").mean()
    rf_monthly.columns = ["RF"]
    return rf_monthly


def fetch_stock_monthly(ticker, start, end):
    """Download monthly ‘Close’ prices and convert to % returns."""
    px = (
        yf.download(
            ticker,
            start=start,
            end=end,
            interval="1mo",
            auto_adjust=False,
            progress=False,
        )["Close"]
        .dropna()
    )
    returns = px.pct_change().mul(100)
    returns.name = "StockRet"
    returns.index = returns.index.to_period("M").to_timestamp("M")
    return returns


def compute_beta(data):
    X = sm.add_constant(data["Mkt-RF"])
    model = sm.OLS(data["Excess"], X).fit()
    return model.params["Mkt-RF"], model


def build_scatter_plot(data, ticker, beta, intercept):
    fig, ax = plt.subplots()
    ax.scatter(data["Mkt-RF"], data["Excess"], alpha=0.7)

    # Regression line
    x_vals = np.linspace(data["Mkt-RF"].min(), data["Mkt-RF"].max(), 100)
    y_vals = intercept + beta * x_vals
    ax.plot(x_vals, y_vals, linewidth=2)

    ax.set_title(f"{ticker} Excess vs. Market Excess Return", fontsize=12)
    ax.set_xlabel("Market Excess Return (Mkt-RF) %")
    ax.set_ylabel(f"{ticker} Excess Return %")
    fig.tight_layout()

    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=300)
    plt.close(fig)
    buf.seek(0)
    return buf


def create_ppt(ticker, plot_bytes, rf_rate, beta, mrp, cost_equity):
    prs = Presentation()

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = f"{ticker} Cost of Equity Capital"
    subtitle = title_slide.placeholders[1]
    subtitle.text = f"Generated {dt.date.today().isoformat()}"

    # Scatter plot slide
    plot_slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    plot_slide.shapes.add_picture(plot_bytes, Inches(0.5), Inches(1), height=Inches(5))

    # Table slide
    table_slide = prs.slides.add_slide(prs.slide_layouts[6])
    rows, cols = 5, 2
    table = table_slide.shapes.add_table(rows, cols, Inches(1), Inches(1.5), Inches(8), Inches(3)).table

    labels = [
        "Metric",
        "Value (%)",
        "Risk‑free rate (RF)",
        f"Beta ({ticker})",
        "Market risk premium (MRP)",
        "Cost of equity (RF + Beta×MRP)",
    ]
    vals = [
        "",
        "",
        f"{rf_rate:.2f}",
        f"{beta:.3f}",
        f"{mrp:.2f}",
        f"{cost_equity:.2f}",
    ]
    for i in range(rows):
        table.cell(i, 0).text = labels[i + 2]
        table.cell(i, 1).text = vals[i + 2]

    # Return PPTX as bytes
    ppt_buf = BytesIO()
    prs.save(ppt_buf)
    ppt_buf.seek(0)
    return ppt_buf


# ---------------------------------------------------------------------------
# Streamlit UI
# ---------------------------------------------------------------------------

st.title("Cost of Equity Capital Deck Generator")
ticker_input = st.text_input("Enter a ticker symbol (e.g., AAPL):")


def run_app(ticker: str):
    # Date range set deliberately long so we have data even for newer firms
    today = dt.date.today()
    start_date = today - dt.timedelta(days=365 * 15)

    # Fetch data
    ff_factors = fetch_market_factors()
    rf_series = fetch_risk_free(start_date, today)
    stock_ret = fetch_stock_monthly(ticker, start_date, today)

    # Combine & clean
    data = pd.concat([stock_ret, ff_factors, rf_series], axis=1, join="inner").dropna()
    data = data.tail(60)  # most recent 60 months
    data["Excess"] = data["StockRet"] - data["RF"]

    # Regression
    beta, model = compute_beta(data)
    intercept = model.params["const"]

    # MRP – annualised long‑run mean of monthly Mkt‑RF
    mrp = ff_factors["Mkt-RF"].mean() * 12

    # Risk‑free rate – latest value in utilised window
    rf_rate = data["RF"].iloc[-1]

    cost_equity = rf_rate + beta * mrp

    # Plot
    plot_bytes = build_scatter_plot(data, ticker, beta, intercept)

    # PowerPoint
    ppt_bytes = create_ppt(ticker, plot_bytes, rf_rate, beta, mrp, cost_equity)

    st.success("PowerPoint deck generated!")
    st.download_button(
        label="Download the deck",
        data=ppt_bytes,
        file_name=f"{ticker}_cost_of_equity.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


if st.button("Generate") and ticker_input:
    try:
        run_app(ticker_input.strip().upper())
    except Exception as exc:
        st.error(f"Error generating deck: {exc}")
